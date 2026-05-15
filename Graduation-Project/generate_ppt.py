#!/tmp/pptx-venv/bin/python3
"""毕业答辩PPT — 最终版"""
from pptx import Presentation, util
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# 只使用四种颜色：黑 蓝 红 绿
BLACK = RGBColor(0x00, 0x00, 0x00)
BLUE  = RGBColor(0x1A, 0x3C, 0x6E)
RED   = RGBColor(0xCC, 0x00, 0x00)
GREEN = RGBColor(0x00, 0x80, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x66, 0x66, 0x66)
LGRAY = RGBColor(0xF2, 0xF2, 0xF2)

def Pt(v): return util.Pt(v)

prs = Presentation()
prs.slide_width  = util.Inches(13.333)
prs.slide_height = util.Inches(7.5)

# ===== 工具函数 =====
def rect(s, l, t, w, h, c, nofill=False):
    x = s.shapes.add_shape(1, util.Inches(l), util.Inches(t), util.Inches(w), util.Inches(h))
    if nofill:
        x.fill.background()
    else:
        x.fill.solid()
        x.fill.fore_color.rgb = c
    x.line.fill.background()
    return x

def txt(s, l, t, w, h, text, sz=16, c=BLACK, b=False, a=PP_ALIGN.LEFT, font="微软雅黑"):
    box = s.shapes.add_textbox(util.Inches(l), util.Inches(t), util.Inches(w), util.Inches(h))
    box.text_frame.word_wrap = True
    box.text_frame.auto_size = None
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.alignment = a
    return box

def bullets(s, l, t, w, h, items, sz=14, c=BLACK, sp=8, prefix="•"):
    box = s.shapes.add_textbox(util.Inches(l), util.Inches(t), util.Inches(w), util.Inches(h))
    box.text_frame.word_wrap = True
    for i, item in enumerate(items):
        p = box.text_frame.paragraphs[0] if i == 0 else box.text_frame.add_paragraph()
        p.text = f"{prefix} {item}"
        p.font.size = Pt(sz)
        p.font.color.rgb = c
        p.space_after = Pt(sp)
    return box

def header(s, title, num=""):
    rect(s, 0, 0, 13.333, 1.1, BLUE)
    if num:
        txt(s, 0.5, 0.15, 1, 0.4, num, 20, WHITE, True)
    txt(s, 1.8 if num else 0.5, 0.2, 10, 0.6, title, 24, WHITE, True)
    rect(s, 0, 1.1, 13.333, 0.04, BLUE)

def section_title(s, l, t, w, text, c=BLUE):
    rect(s, l, t, 0.08, 0.5, c)
    txt(s, l+0.25, t+0.05, w-0.5, 0.4, text, 18, BLACK, True)

# ============================================================
# P1 封面（黑色改为蓝色）
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, 13.333, 7.5, BLUE)
rect(s, 0, 0, 13.333, 0.06, WHITE)
rect(s, 0, 7.44, 13.333, 0.06, WHITE)

txt(s, 1.5, 1.0, 10, 0.6, "本科毕业设计答辩", 16, WHITE)
txt(s, 1.5, 2.0, 10, 1.2, "基于优化XGBoost的DGA检测系统设计", 34, WHITE, True)
rect(s, 1.5, 3.3, 3, 0.04, WHITE)

info = [
    ("学    院", "数学与信息科学学院"),
    ("专    业", "信息安全"),
    ("班    级", "信安221"),
    ("学生姓名", "苑文洋"),
    ("学    号", "32215300032"),
    ("指导教师", "余玉银"),
]
for i, (k, v) in enumerate(info):
    txt(s, 1.5, 3.6+i*0.45, 2, 0.35, k, 14, WHITE)
    txt(s, 3.6, 3.6+i*0.45, 5, 0.35, v, 14, WHITE)

# ============================================================
# P2-3 前言（4部分放2页）
# ============================================================
# P2 前言（上）
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "前  言", "01")

section_title(s, 0.5, 1.5, 12, "1. 研究背景", BLUE)
bullets(s, 0.8, 2.1, 11.5, 2.0, [
    "DGA（Domain Generation Algorithm）是僵尸网络逃避黑名单检测的核心手段",
    "攻击者通过算法生成大量随机域名，动态切换C&C服务器",
    "传统黑名单、签名检测等方法难以应对域名频繁变种",
    "每日可生成数万候选域名，防御方无法提前预判封锁",
], 14, BLACK, 8)

section_title(s, 0.5, 4.0, 12, "2. 研究意义", BLUE)
bullets(s, 0.8, 4.6, 11.5, 2.0, [
    "DGA检测是网络安全主动防御的重要环节",
    "高效的DGA检测可阻断僵尸网络C2通信链路",
    "机器学习方法具备泛化能力，可检测未知变种域名",
    "XGBoost算法在效率与精度间取得良好平衡",
], 14, BLACK, 8)

# P3 前言（下）
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "前  言", "02")

section_title(s, 0.5, 1.5, 12, "3. 国内外研究现状", BLUE)
bullets(s, 0.8, 2.1, 11.5, 2.5, [
    "传统方法：基于黑名单匹配、规则签名、DNS流量分析",
    "深度学习方法：CNN/LSTM等模型检测精度高，但计算开销大",
    "集成学习方法：随机森林、XGBoost等兼顾效率与可解释性",
    "现有研究多关注二分类，对多分类家族溯源研究较少",
], 14, BLACK, 8)

section_title(s, 0.5, 4.5, 12, "4. 本文主要工作", BLUE)
bullets(s, 0.8, 5.1, 11.5, 1.8, [
    "构建基于XGBoost的DGA检测系统，支持二分类与多分类",
    "对比分析统计特征、N-gram特征及融合特征方案",
    "引入Optuna贝叶斯优化实现超参数自动寻优",
    "二分类F1达98.40%，多分类宏平均F1达88.95%",
], 14, BLACK, 8)

# ============================================================
# P4 DGA域名
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "相关理论与技术——DGA域名", "03")

section_title(s, 0.5, 1.5, 12, "DGA域名工作原理", BLUE)
bullets(s, 0.8, 2.1, 11, 2.5, [
    "DGA（Domain Generation Algorithm）：恶意软件用于动态生成域名的算法",
    "种子（种子密钥）+ 算法逻辑 → 每日/每周生成成百上千候选域名",
    "攻击者只需注册其中少数几个，即可维持C2（命令与控制）通信",
    "安全团队难以预判和封锁所有生成域名，形成检测盲区",
], 14, BLACK, 8)

section_title(s, 0.5, 4.5, 12, "常见DGA家族", BLUE)
bullets(s, 0.8, 5.1, 11, 1.8, [
    "Suppobox：基于日期和字典的DGA，生成可读性较高的域名",
    "Kraken：字符拼接型DGA，域名长度变化较大",
    "Shiotob：基于时间种子的DGA，字符随机性强",
    "Bamital：固定结构的DGA，域名模式相对稳定",
], 14, BLACK, 8)

# ============================================================
# P5 XGBoost
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "相关理论与技术——XGBoost", "04")

section_title(s, 0.5, 1.5, 12, "XGBoost核心原理", BLUE)
bullets(s, 0.8, 2.1, 11, 2.0, [
    "XGBoost（eXtreme Gradient Boosting）：基于梯度提升决策树的优化算法",
    "逐步添加决策树，每棵新树都在修正前一棵树的残差（Boosting思想）",
    "目标函数 = 训练损失 + 正则化项（L1+L2），防止过拟合",
    "使用二阶泰勒展开近似损失函数，收敛速度比传统GBDT更快",
], 14, BLACK, 8)

section_title(s, 0.5, 4.0, 12, "选择XGBoost的原因", BLUE)
bullets(s, 0.8, 4.6, 11, 2.0, [
    "天然适合高维稀疏特征（N-gram特征10,000维）",
    "相比深度学习方法训练速度快5-10倍，可解释性强",
    "内置列采样、缩减（Shrinkage）、早停等防过拟合机制",
    "特征重要性可直接输出，便于分析各维度贡献",
], 14, BLACK, 8)

# ============================================================
# P6 贝叶斯调参
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "相关理论与技术——贝叶斯优化", "05")

section_title(s, 0.5, 1.5, 12, "贝叶斯优化原理", BLUE)
bullets(s, 0.8, 2.1, 11, 2.0, [
    "基于高斯过程的超参数优化方法，利用历史评估结果指导下一步搜索",
    "相比网格搜索（GridSearch）：搜索效率提升10倍以上",
    "相比随机搜索（RandomSearch）：能更聚焦于最优区域",
    "Optuna框架：支持自动剪枝（Pruning）、可视化、分布式并行",
], 14, BLACK, 8)

section_title(s, 0.5, 4.0, 12, "超参数搜索范围", BLUE)
# 参数表
tbl = s.shapes.add_table(7, 2, util.Inches(1), util.Inches(4.6), util.Inches(6), util.Inches(2.5)).table
tbl.cell(0,0).text = "超参数"; tbl.cell(0,1).text = "搜索范围"
for c in [tbl.cell(0,0), tbl.cell(0,1)]:
    for p in c.text_frame.paragraphs: p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = WHITE
    c.fill.solid(); c.fill.fore_color.rgb = BLUE

params = [("learning_rate", "[0.01, 0.3]"),("max_depth", "[3, 12]"),
          ("n_estimators", "[50, 500]"),("subsample", "[0.6, 1.0]"),
          ("colsample_bytree", "[0.6, 1.0]"),("reg_lambda", "[0, 5]")]
for i, (k, v) in enumerate(params):
    tbl.cell(i+1,0).text = k; tbl.cell(i+1,1).text = v
    for j in range(2):
        for p in tbl.cell(i+1,j).text_frame.paragraphs: p.font.size = Pt(11)
    if i % 2 == 0:
        for j in range(2): tbl.cell(i+1,j).fill.solid(); tbl.cell(i+1,j).fill.fore_color.rgb = LGRAY

# ============================================================
# P7 评价指标
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "相关理论与技术——评价指标", "06")

section_title(s, 0.5, 1.5, 12, "模型评估指标", BLUE)
bullets(s, 0.8, 2.1, 5.5, 3.5, [
    "准确率（Accuracy）：正确分类样本占总样本比例",
    "精确率（Precision）：预测为恶意的样本中真正恶意的比例",
    "召回率（Recall）：所有恶意样本中被正确识别的比例",
    "F1-score：精确率和召回率的调和平均数",
    "AUC-ROC：ROC曲线下面积，衡量模型整体区分能力",
], 14, BLACK, 8)

# 公式区域
rect(s, 7, 2.1, 5.5, 3.5, LGRAY)
txt(s, 7.3, 2.2, 5, 0.4, "核心公式", 16, BLUE, True)
txt(s, 7.3, 2.8, 5, 0.6, "F1 = 2 · Precision · Recall / (Precision + Recall)", 13, BLACK)
txt(s, 7.3, 3.5, 5, 0.6, "Accuracy = (TP + TN) / (TP + TN + FP + FN)", 13, BLACK)
txt(s, 7.3, 4.2, 5, 0.6, "Precision = TP / (TP + FP)", 13, BLACK)
txt(s, 7.3, 4.9, 5, 0.5, "Recall = TP / (TP + FN)", 13, BLACK)

# ============================================================
# P8-P10 DGA域名检测方法设计（1,2,3各一页）
# ============================================================
# P8 方法设计1：数据采集与预处理
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "DGA域名检测方法设计——数据采集与预处理", "07")

section_title(s, 0.5, 1.5, 12, "1. 数据采集", BLUE)
bullets(s, 0.8, 2.1, 11, 1.5, [
    "良性域名来源：Alexa全球流量榜单，选取20万条正常域名",
    "恶意域名来源：360 DGA监控库，选取20万条DGA恶意域名",
    "合计40万条样本，涵盖多个主流DGA家族",
], 14, BLACK, 8)

section_title(s, 0.5, 3.8, 12, "2. 数据预处理", BLUE)
bullets(s, 0.8, 4.4, 11, 2.5, [
    "数据清洗：去除空值、重复域名，保证数据质量",
    "TLD剥离：去除.com/.net/.org等顶级域后缀，聚焦二级域名特征",
    "标签编码：良性域名标为0，各DGA家族标为1-10",
    "数据集划分：训练集80%（32万）、验证集10%（4万）、测试集10%（4万）",
], 14, BLACK, 8)

# P9 方法设计2：特征提取
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "DGA域名检测方法设计——特征提取", "08")

section_title(s, 0.5, 1.5, 12, "3. 特征提取", BLUE)
bullets(s, 0.8, 2.1, 5.5, 3.0, [
    "统计特征（6维）：",
    "  域名长度：正常与DGA域名字符数差异",
    "  数字占比：DGA域名常含大量随机数字",
    "  信息熵：随机字符串的熵值更高",
    "  连续辅音比例：恶意域名常出现辅音堆叠",
    "  重复字符比例、元音占比",
], 13, BLACK, 6)

bullets(s, 7, 2.1, 5.5, 3.0, [
    "N-gram特征：",
    "  将域名按n个连续字符滑动切分",
    "  2-gram：apple → ap pp pl le",
    "  2-4gram：混合2/3/4字符组合",
    "  CountVectorizer统计各组合频次",
    "  取前10,000个高频组合作为特征维度",
], 13, BLACK, 6)

rect(s, 0.5, 5.5, 12.3, 1.2, LGRAY)
txt(s, 0.8, 5.6, 11, 0.4, "💡 融合特征策略", 14, BLUE, True)
txt(s, 0.8, 6.0, 11, 0.5, "统计特征（宏观基准定位）+ 2-gram特征（微观指纹修正）→ 维度 6 + 10,000 = 10,006", 13, BLACK)

# P10 方法设计3：模型训练
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "DGA域名检测方法设计——模型训练", "09")

section_title(s, 0.5, 1.5, 12, "4. 模型训练", BLUE)

# 对比表
tbl = s.shapes.add_table(5, 4, util.Inches(0.8), util.Inches(2.2), util.Inches(11.5), util.Inches(2.5)).table
headers = ["特征组", "特征组合", "说明", "维度"]
for j, h in enumerate(headers):
    tbl.cell(0,j).text = h
    for p in tbl.cell(0,j).text_frame.paragraphs: p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = WHITE
    tbl.cell(0,j).fill.solid(); tbl.cell(0,j).fill.fore_color.rgb = BLUE
data = [["①", "统计特征", "6维统计特征", "6"],
        ["②", "2-gram 特征", "Bigram + CountVectorizer", "10,000"],
        ["③", "2-4gram 特征", "混合切片", "10,000"],
        ["④", "统计 + 2-gram", "融合特征", "10,006"]]
for i, row in enumerate(data):
    for j, val in enumerate(row):
        tbl.cell(i+1,j).text = val
        for p in tbl.cell(i+1,j).text_frame.paragraphs: p.font.size = Pt(11)
        if i % 2 == 0: tbl.cell(i+1,j).fill.solid(); tbl.cell(i+1,j).fill.fore_color.rgb = LGRAY

bullets(s, 0.8, 5.0, 11, 1.5, [
    "每组特征分别在 XGBoost 默认参数和 Optuna 贝叶斯优化下训练",
    "二分类任务：区分恶意域名（DGA）与良性域名（Alexa）",
    "多分类任务：对10个DGA家族进行细粒度分类溯源",
], 13, BLACK, 6)

# ============================================================
# P11-P12 二分类（2页）
# ============================================================
# P11 二分类1
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "实验与结果分析——二分类结果（上）", "10")

rect(s, 0.5, 1.4, 12.3, 0.8, LGRAY)
txt(s, 0.8, 1.5, 11, 0.4, "🏆 最优结果：统计特征 + 2-gram 融合组 + Optuna 优化", 16, RED, True)
txt(s, 0.8, 1.85, 11, 0.3, "准确率 98.40%  |  F1 98.40%  |  AUC 0.9989", 14, GREEN, True)

tbl = s.shapes.add_table(5, 7, util.Inches(0.3), util.Inches(2.6), util.Inches(12.7), util.Inches(2.8)).table
h2 = ["特征组", "优化", "准确率", "精确率", "召回率", "F1", "AUC"]
for j, h in enumerate(h2):
    tbl.cell(0,j).text = h
    for p in tbl.cell(0,j).text_frame.paragraphs: p.font.size = Pt(10); p.font.bold = True; p.font.color.rgb = WHITE
    tbl.cell(0,j).fill.solid(); tbl.cell(0,j).fill.fore_color.rgb = BLUE
d2 = [["统计特征", "默认", "95.87%", "95.97%", "95.77%", "95.87%", "0.9931"],
      ["统计特征", "Optuna","96.12%","96.48%","95.76%","96.12%","0.9948"],
      ["2-gram", "默认", "97.53%", "97.44%", "97.63%", "97.53%", "0.9979"],
      ["2-gram", "Optuna","97.69%","97.72%","97.66%","97.69%","0.9981"]]
for i, row in enumerate(d2):
    for j, val in enumerate(row):
        tbl.cell(i+1,j).text = val
        for p in tbl.cell(i+1,j).text_frame.paragraphs: p.font.size = Pt(10)
    if i % 2 == 0:
        for j in range(7): tbl.cell(i+1,j).fill.solid(); tbl.cell(i+1,j).fill.fore_color.rgb = LGRAY

txt(s, 0.5, 5.8, 11, 0.4, "📌 分析：统计特征单独使用时，准确率约95.87%；2-gram特征表现更好，达到97.53%", 13, BLACK)
txt(s, 0.5, 6.2, 11, 0.3, "Optuna优化在两组特征上均有小幅提升（约+0.2~0.3%），说明默认参数已接近最优", 13, GRAY)

# P12 二分类2
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "实验与结果分析——二分类结果（下）", "11")

tbl = s.shapes.add_table(3, 7, util.Inches(0.5), util.Inches(1.6), util.Inches(12.3), util.Inches(1.8)).table
for j, h in enumerate(h2):
    tbl.cell(0,j).text = h
    for p in tbl.cell(0,j).text_frame.paragraphs: p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = WHITE
    tbl.cell(0,j).fill.solid(); tbl.cell(0,j).fill.fore_color.rgb = BLUE
d2b = [["统计+2-gram", "默认", "98.10%", "98.39%", "97.81%", "98.10%", "0.9984"],
       ["统计+2-gram", "Optuna","98.40%","98.47%","98.33%","98.40%","0.9989"]]
for i, row in enumerate(d2b):
    for j, val in enumerate(row):
        tbl.cell(i+1,j).text = val
        for p in tbl.cell(i+1,j).text_frame.paragraphs: p.font.size = Pt(11); p.font.color.rgb = RED if i==1 else BLACK
        if i == 1: p.font.bold = True
    if i == 1:
        for j in range(7): tbl.cell(i+1,j).fill.solid(); tbl.cell(i+1,j).fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xF0)

bullets(s, 0.5, 3.8, 11.5, 3.0, [
    "融合特征（统计+2-gram）效果显著优于单一特征",
    "默认参数下融合组准确率98.10%，高于单一特征的97.69%和96.12%",
    "Optuna优化后融合组达98.40%，F1 98.40%，AUC 0.9989",
    "融合特征实现了统计特征（宏观）与2-gram特征（微观）的互补增益"
], 14, BLACK, 8)

# ============================================================
# P13-P14 多分类（2页）
# ============================================================
# P13 多分类1
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "实验与结果分析——多分类结果（上）", "12")

tbl = s.shapes.add_table(5, 4, util.Inches(1), util.Inches(1.6), util.Inches(11), util.Inches(2.2)).table
for j, h in enumerate(["特征组", "准确率", "宏平均F1", "说明"]):
    tbl.cell(0,j).text = h
    for p in tbl.cell(0,j).text_frame.paragraphs: p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = WHITE
    tbl.cell(0,j).fill.solid(); tbl.cell(0,j).fill.fore_color.rgb = BLUE
d3 = [["统计特征", "65.59%", "53.78%", "宏观基准定位，区分度有限"],
      ["2-gram 特征", "84.46%", "76.80%", "微观指纹提取，效果显著"],
      ["2-4gram 特征", "85.55%", "77.70%", "相比2-gram提升有限"],
      ["统计+2-gram", "89.98%", "88.95%", "✅ 最优组合"]]
for i, row in enumerate(d3):
    for j, val in enumerate(row):
        tbl.cell(i+1,j).text = val
        for p in tbl.cell(i+1,j).text_frame.paragraphs: p.font.size = Pt(11)
        if i == 3: p.font.color.rgb = GREEN; p.font.bold = True
    if i % 2 == 0:
        for j in range(4): tbl.cell(i+1,j).fill.solid(); tbl.cell(i+1,j).fill.fore_color.rgb = LGRAY

section_title(s, 0.5, 4.2, 12, "分析发现", BLUE)
bullets(s, 0.8, 4.8, 11.5, 2.0, [
    "统计特征单独用于多分类效果有限（F1仅53.78%），因不同家族统计特性差异不大",
    "2-gram特征F1提升至76.80%，说明字符序列模式能有效区分家族",
    "2-4gram相比2-gram仅提升0.9%，说明增加n值边际收益递减",
    "融合特征F1达88.95%，验证了宏观+微观的互补有效性",
], 13, BLACK, 7)

# P14 多分类2
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "实验与结果分析——多分类结果（下）", "13")

section_title(s, 0.5, 1.5, 12, "各DGA家族细粒度F1对比", BLUE)

tbl = s.shapes.add_table(6, 4, util.Inches(0.8), util.Inches(2.1), util.Inches(11), util.Inches(2.8)).table
for j, h in enumerate(["DGA家族", "F1-score", "样本量", "特征分析"]):
    tbl.cell(0,j).text = h
    for p in tbl.cell(0,j).text_frame.paragraphs: p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = WHITE
    tbl.cell(0,j).fill.solid(); tbl.cell(0,j).fill.fore_color.rgb = BLUE
f1_data = [["Kraken", "92.3%", "40,000", "长度+字符分布特征明显"],
           ["Shiotob", "91.8%", "35,000", "n-gram模式区分度高"],
           ["Bamital", "90.5%", "30,000", "域名结构较为固定"],
           ["Pykspa", "87.2%", "25,000", "随机化程度较高"],
           ["Suppobox", "72.1%", "8,000", "样本量有限，模式随机"]]
for i, row in enumerate(f1_data):
    for j, val in enumerate(row):
        tbl.cell(i+1,j).text = val
        for p in tbl.cell(i+1,j).text_frame.paragraphs: p.font.size = Pt(11)
    if i % 2 == 0:
        for j in range(4): tbl.cell(i+1,j).fill.solid(); tbl.cell(i+1,j).fill.fore_color.rgb = LGRAY

bullets(s, 0.8, 5.3, 11.5, 1.5, [
    "样本量超过25,000的家族F1均达87%以上，模型表现稳定",
    "Suppobox家族F1偏低（72.1%），主要受限于样本量不足（仅8,000条）",
    "样本量对多分类精度有显著影响，增加训练数据可进一步提升性能",
], 13, BLACK, 7)

# ============================================================
# P15 小结
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "小  结", "14")

section_title(s, 0.5, 1.5, 12, "本文工作总结", BLUE)
bullets(s, 0.8, 2.1, 11, 2.5, [
    "构建了基于XGBoost的DGA检测系统，同时支持二分类（恶意检测）和多分类（家族溯源）",
    '系统对比了统计特征、N-gram特征及融合特征四种方案，"统计+2-gram"为最优',
    "引入Optuna贝叶斯优化自动调参，相比默认参数有稳定提升",
    "二分类任务最优F1达98.40%，AUC达0.9989",
    "多分类任务宏平均F1达88.95%",
], 14, BLACK, 8)

section_title(s, 0.5, 4.8, 12, "本文创新点", BLUE)
bullets(s, 0.8, 5.4, 11, 1.5, [
    "🌟 特征融合策略：揭示统计特征（宏观）与2-gram特征（微观）的正交互补机制",
    "🌟 优化方法：引入Optuna替代网格搜索，调参效率提升10倍以上",
    "🌟 统一框架：同时支持恶意域名检测和DGA家族溯源",
], 14, BLACK, 8)

# ============================================================
# P16 未来展望
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "未来展望", "15")

section_title(s, 0.5, 1.5, 12, "当前不足", BLUE)
bullets(s, 0.8, 2.1, 11, 1.5, [
    "多分类准确率仍有提升空间，部分DGA家族（如Suppobox）样本量不足",
    "目前仅在离线数据集上进行验证，未部署到在线实时检测环境",
    "模型在极端不平衡样本下的表现有待进一步验证",
], 13, BLACK, 8)

section_title(s, 0.5, 3.8, 12, "未来改进方向", BLUE)
bullets(s, 0.8, 4.4, 11, 2.0, [
    "引入深度学习模型（CNN/LSTM）与XGBoost集成，构建混合检测框架",
    "部署在线实时检测系统，在真实网络流量中验证模型效果",
    "增加更多DGA家族样本覆盖，提升多分类的粒度与精度",
    "探索基于主动学习的半监督检测方法，降低对标注数据的依赖",
], 14, BLACK, 8)

# ============================================================
# P17 结束页
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, 13.333, 7.5, BLUE)
rect(s, 0, 0, 13.333, 0.06, WHITE)
rect(s, 0, 7.44, 13.333, 0.06, WHITE)

txt(s, 2, 1.5, 9, 1.0, "致  谢", 44, WHITE, True, PP_ALIGN.CENTER)
rect(s, 5, 2.6, 3.3, 0.03, WHITE)

txt(s, 2, 3.0, 9, 1.5, "感谢指导老师余玉银老师的悉心指导与耐心帮助\n"
         "感谢数学与信息科学学院各位老师的教诲与培养\n"
         "感谢家人和同学的支持与鼓励",
         18, WHITE, a=PP_ALIGN.CENTER)

rect(s, 4.5, 5.0, 4.3, 0.03, WHITE)
txt(s, 2, 5.3, 9, 1.0, "请各位老师批评指正！", 28, WHITE, True, PP_ALIGN.CENTER)
txt(s, 2, 6.3, 9, 0.5, "Q & A", 18, WHITE, a=PP_ALIGN.CENTER)

# ===== 保存 =====
out = "/home/girlorn/Cloud-Security-Study/Graduation-Project/毕业答辩PPT.pptx"
prs.save(out)
print(f"✅ 保存: {out}")
print(f"共 {len(prs.slides)} 页")
